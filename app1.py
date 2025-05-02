import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from gtts import gTTS
import base64
from PIL import Image, ImageOps, ImageFilter
import pytesseract
from pptx import Presentation
import ffmpeg
import tempfile
from pydantic import Field, ConfigDict
import cv2
import numpy as np
import io
from skimage.transform import rotate
from deskew import determine_skew
import matplotlib.pyplot as plt
import speech_recognition as sr
import pyttsx3
from datetime import datetime, timedelta
import hashlib
import sqlite3
import re
from passlib.hash import bcrypt, sha256_crypt
from jose import jwt
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import secrets
import string
import threading
import time

# Load environment variables from .env
load_dotenv()

# Set your Google API key here - Please replace with your own API key in production
os.environ["GOOGLE_API_KEY"] = os.getenv("GOOGLE_API_KEY")



# Initialize text-to-speech engine
tts_engine = pyttsx3.init()
tts_engine.setProperty('rate', 150)  # Adjust speaking rate
tts_engine.setProperty('volume', 0.9)  # Adjust volume
tts_engine.setProperty('voice', 'english-us')  # Set voice to US English
tts_engine.setProperty('pitch', 50)  # Adjust pitch

# Initialize database
def init_db():
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS users
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  username TEXT UNIQUE,
                  password TEXT,
                  email TEXT UNIQUE,
                  first_name TEXT,
                  last_name TEXT,
                  role TEXT,
                  created_at TIMESTAMP)''')
    conn.commit()
    conn.close()

init_db()

# Updated hash_password function
def hash_password(password: str) -> str:
    """Hash password using bcrypt (for new users)."""
    return bcrypt.hash(password)

# Updated verify_password function
def verify_password(plain_password: str, hashed_password: str) -> bool:
    """Verify password with automatic hash migration."""
    # First try bcrypt
    if hashed_password.startswith("$2b$"):
        return bcrypt.verify(plain_password, hashed_password)
    
    # Fallback to SHA-256 (old method)
    if hashed_password == hashlib.sha256(plain_password.encode()).hexdigest():
        # Rehash with bcrypt if using old method
        return True
    return False

# Email validation
def is_valid_email(email):
    pattern = r"^[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+$"
    return re.match(pattern, email)

# Add email verification function
def send_verification_email(email, token):
    """Send a verification email with a token."""
    try:
        smtp_server = "smtp.example.com"  # Replace with your SMTP server
        smtp_port = 587  # Replace with your SMTP port
        smtp_user = "your_email@example.com"  # Replace with your email
        smtp_password = "your_password"  # Replace with your email password

        subject = "Verify Your Email"
        body = f"Please verify your email by clicking the link: https://example.com/verify?token={token}"
        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = smtp_user
        msg["To"] = email

        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(smtp_user, smtp_password)
            server.sendmail(smtp_user, email, msg.as_string())
        return True
    except Exception as e:
        st.error(f"Failed to send verification email: {str(e)}")
        return False

# Email configuration
EMAIL_HOST = os.getenv("EMAIL_HOST", "smtp.gmail.com")
EMAIL_PORT = os.getenv("EMAIL_PORT", 587)
EMAIL_USER = os.getenv("EMAIL_USER", "your-email@gmail.com")
EMAIL_PASS = os.getenv("EMAIL_PASS", "your-app-password")
BASE_URL = os.getenv("BASE_URL", "http://localhost:8501")

def generate_reset_token():
    """Generate a secure random token."""
    return ''.join(secrets.choice(string.ascii_letters + string.digits) for _ in range(32))

def send_reset_email(email: str, token: str):
    """Send password reset email."""
    try:
        msg = MIMEMultipart()
        msg['From'] = EMAIL_USER
        msg['To'] = email
        msg['Subject'] = "MediCare-ICU Password Reset"
        
        reset_link = f"{BASE_URL}/?token={token}"
        body = f"""
        <h2>Password Reset Request</h2>
        <p>Click the link below to reset your password:</p>
        <p><a href="{reset_link}">Reset Password</a></p>
        <p>This link will expire in 1 hour.</p>
        <p>If you didn't request this, please ignore this email.</p>
        """
        msg.attach(MIMEText(body, 'html'))
        
        with smtplib.SMTP(EMAIL_HOST, EMAIL_PORT) as server:
            server.starttls()
            server.login(EMAIL_USER, EMAIL_PASS)
            server.send_message(msg)
        return True
    except Exception as e:
        st.error(f"Error sending email: {str(e)}")
        return False

def store_reset_token(email: str, token: str):
    """Store token in database with expiration."""
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    expires_at = datetime.now() + timedelta(hours=1)
    try:
        c.execute("""
            INSERT OR REPLACE INTO password_resets 
            (email, token, expires_at) 
            VALUES (?, ?, ?)
        """, (email, token, expires_at))
        conn.commit()
        return True
    except Exception as e:
        st.error(f"Database error: {str(e)}")
        return False
    finally:
        conn.close()

def show_password_reset():
    """Password reset UI."""
    with st.form("reset_form"):
        email = st.text_input("Enter your registered email")
        if st.form_submit_button("Send Reset Link"):
            conn = sqlite3.connect('medicare_users.db')
            c = conn.cursor()
            c.execute("SELECT email FROM users WHERE email=?", (email,))
            if not c.fetchone():
                st.error("Email not found in our system")
                return
            
            token = generate_reset_token()
            if store_reset_token(email, token) and send_reset_email(email, token):
                st.success("Reset link sent to your email!")
            else:
                st.error("Failed to send reset email")

def handle_password_reset():
    """Check for reset token in URL."""
    query_params = st.query_params  # Updated from st.experimental_get_query_params
    if 'token' in query_params:
        token = query_params['token']
        conn = sqlite3.connect('medicare_users.db')
        c = conn.cursor()
        c.execute("""
            SELECT email FROM password_resets 
            WHERE token=? AND expires_at > datetime('now')
        """, (token,))
        result = c.fetchone()
        
        if result:
            email = result[0]
            with st.form("new_password_form"):
                st.subheader("Set New Password")
                new_password = st.text_input("New Password", type="password")
                confirm_password = st.text_input("Confirm Password", type="password")
                
                if st.form_submit_button("Update Password"):
                    if new_password != confirm_password:
                        st.error("Passwords don't match")
                    elif len(new_password) < 8:
                        st.error("Password must be at least 8 characters")
                    else:
                        hashed = hash_password(new_password)
                        c.execute("UPDATE users SET password=? WHERE email=?", (hashed, email))
                        c.execute("DELETE FROM password_resets WHERE email=?", (email,))
                        conn.commit()
                        st.success("Password updated successfully! Please login.")
                        st.query_params.clear()  # Updated from st.experimental_set_query_params
                        st.rerun()
        else:
            st.error("Invalid or expired reset link")
        conn.close()

# Simple Authentication System
class AuthSystem:
    def __init__(self):
        self._user = None
    
    @property
    def user(self):
        if 'user' not in st.session_state:
            return None
        return st.session_state.user
    
    def sign_up(self, username, password, email, first_name, last_name, role="patient"):
        if not is_valid_email(email):
            return False, "Invalid email format"
        
        if len(password) < 8:
            return False, "Password must be at least 8 characters"
        
        conn = sqlite3.connect('medicare_users.db')
        c = conn.cursor()
        try:
            hashed_password = hash_password(password)
            c.execute("INSERT INTO users (username, password, email, first_name, last_name, role, created_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
                     (username, hashed_password, email, first_name, last_name, role, datetime.now()))
            conn.commit()

            # Generate a token for email verification
            token = jwt.encode({"email": email}, "your_secret_key", algorithm="HS256")
            if send_verification_email(email, token):
                return True, "Account created successfully! Please verify your email."
            else:
                return False, "Account created, but failed to send verification email."
        except sqlite3.IntegrityError as e:
            if "username" in str(e):
                return False, "Username already exists"
            elif "email" in str(e):
                return False, "Email already registered"
            return False, "Registration failed"
        finally:
            conn.close()
    
    def sign_in(self, username, password):
        conn = sqlite3.connect('medicare_users.db')
        c = conn.cursor()
        c.execute("SELECT * FROM users WHERE username=?", (username,))
        user = c.fetchone()
        
        if user:
            stored_hash = user[2]
            if verify_password(password, stored_hash):
                # If using old hash, upgrade to bcrypt
                if not stored_hash.startswith("$2b$"):
                    new_hash = hash_password(password)
                    c.execute("UPDATE users SET password=? WHERE username=?", 
                             (new_hash, username))
                    conn.commit()
                
                # Set user session
                st.session_state.user = {
                    "id": user[0],
                    "username": user[1],
                    "email": user[3],
                    "firstName": user[4],
                    "lastName": user[5],
                    "role": user[6]
                }
                conn.close()
                return True, "Logged in successfully!"
        
        conn.close()
        return False, "Invalid username or password"
    
    def sign_out(self):
        if 'user' in st.session_state:
            del st.session_state.user
        st.rerun()

auth = AuthSystem()

# One-time password reset script for migration
def migrate_hashes():
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    c.execute("SELECT id, username, password FROM users")
    users = c.fetchall()
    
    for user in users:
        user_id, username, old_hash = user
        # Skip already migrated users
        if old_hash.startswith("$2b$"):
            continue
            
        # Manually reset password (users will need to use "forgot password")
        c.execute("UPDATE users SET password=NULL WHERE id=?", (user_id,))
    
    conn.commit()
    conn.close()
    print("Password migration complete. Users must reset passwords.")

# Password reset functionality
def send_password_reset(email):
    """Send password reset email."""
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    c.execute("SELECT id FROM users WHERE email=?", (email,))
    user = c.fetchone()
    
    if user:
        # Generate token (in production, use proper JWT)
        token = hashlib.sha256(f"{email}{datetime.now()}".encode()).hexdigest()
        # Store token in DB (implementation omitted)
        # Send email (implementation omitted)
        return True
    return False

def show_password_reset():
    """Password reset UI."""
    with st.form("reset_form"):
        email = st.text_input("Registered Email")
        if st.form_submit_button("Send Reset Link"):
            if send_password_reset(email):
                st.success("Reset link sent to your email")
            else:
                st.error("Email not found")

# Updated authentication UI
def show_auth():
    tabs = st.tabs(["Login", "Sign Up", "Forgot Password?"])
    
    with tabs[0]:  # Login
        with st.form("login_form"):
            st.subheader("Login to MediCare-ICU")
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            
            if st.form_submit_button("Login"):
                success, message = auth.sign_in(username, password)
                if success:
                    st.success(message)
                    st.rerun()
                else:
                    st.error(message)
    
    with tabs[1]:  # Sign Up
        with st.form("signup_form"):
            st.subheader("Create New Account")
            col1, col2 = st.columns(2)
            with col1:
                first_name = st.text_input("First Name*", placeholder="John")
            with col2:
                last_name = st.text_input("Last Name*", placeholder="Doe")
            
            username = st.text_input("Username*", placeholder="johndoe123")
            email = st.text_input("Email*", placeholder="john@example.com")
            password = st.text_input("Password* (min 8 chars)", type="password")
            confirm_password = st.text_input("Confirm Password*", type="password")
            role = st.selectbox("Role*", ["Patient", "Doctor", "Nurse", "Researcher"])
            
            if st.form_submit_button("Create Account"):
                if password != confirm_password:
                    st.error("Passwords don't match!")
                else:
                    success, message = auth.sign_up(username, password, email, first_name, last_name, role)
                    if success:
                        st.success(message)
                        st.session_state.auth_tab = "login"
                        st.rerun()
                    else:
                        st.error(message)
    
    with tabs[2]:  # Password Reset
        show_password_reset()

def check_authentication():
    """Check if user is authenticated"""
    if not auth.user:
        with st.expander("🔐 Sign In Required", expanded=True):
            st.warning("Please sign in to access the MediCare-ICU Assistant")
            
            # Simple demo authentication
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            
            if st.button("Sign In"):
                # Replace with your actual authentication logic
                if username == "demo" and password == "medicare123":
                    st.session_state.user = {
                        "firstName": "Demo",
                        "lastName": "User",
                        "email": "demo@medicare.com"
                    }
                    st.success("Logged in successfully!")
                    st.rerun()
                else:
                    st.error("Invalid credentials")
        st.stop()

def generate_audio_file(text: str, filename: str) -> str:
    """Generate an audio file from text using gTTS and save it locally."""
    try:
        tts = gTTS(text=text, lang='en')
        tts.save(filename)
        return filename
    except Exception as e:
        st.error(f"Error generating audio file: {str(e)}")
        return None

def autoplay_audio(file_path: str):
    """Generate HTML for autoplaying a saved audio file."""
    try:
        with open(file_path, "rb") as f:
            data = f.read()
            b64 = base64.b64encode(data).decode()
            audio_html = f"""
                <div class="audio-container">
                    <audio controls autoplay>
                        <source src="data:audio/mp3;base64,{b64}" type="audio/mp3">
                        Your browser does not support the audio element.
                    </audio>
                </div>
            """
            st.markdown(audio_html, unsafe_allow_html=True)
    except Exception as e:
        st.error(f"An error occurred while playing audio: {str(e)}")

# Image Processing Functions
def correct_orientation(image):
    """Correct orientation issues from laptop camera"""
    try:
        # Convert PIL Image to numpy array if needed
        if isinstance(image, Image.Image):
            image_np = np.array(image)
        else:
            image_np = image.copy()
        
        return Image.fromarray(image_np)
    except Exception as e:
        st.error(f"Orientation correction failed: {str(e)}")
        return Image.fromarray(image_np) if 'image_np' in locals() else image

def enhance_blurry_image(image, sharpen_factor=2.0, denoise_strength=10):
    """Enhance blurry document images for better OCR"""
    try:
        # Convert to numpy array if PIL Image
        if isinstance(image, Image.Image):
            img = np.array(image)
        else:
            img = image.copy()
        
        # Convert to grayscale
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
        else:
            gray = img
        
        # Apply adaptive histogram equalization
        clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
        enhanced = clahe.apply(gray)
        
        # Denoising
        denoised = cv2.fastNlMeansDenoising(enhanced, None, denoise_strength, 7, 21)
        
        # Sharpening
        kernel = np.array([[-1,-1,-1], [-1,9,-1], [-1,-1,-1]]) * sharpen_factor
        sharpened = cv2.filter2D(denoised, -1, kernel)
        
        # Adaptive thresholding
        thresholded = cv2.adaptiveThreshold(sharpened, 255, 
                                          cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                          cv2.THRESH_BINARY, 11, 2)
        
        return Image.fromarray(thresholded)
    except Exception as e:
        st.error(f"Image enhancement failed: {str(e)}")
        return image

def preprocess_image(image):
    """Enhance document image for better OCR with blur detection"""
    try:
        # Convert to numpy array if PIL Image
        if isinstance(image, Image.Image):
            img = np.array(image)
        else:
            img = image.copy()
        
        # Check if image is blurry
        gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY) if len(img.shape) == 3 else img
        blur_value = cv2.Laplacian(gray, cv2.CV_64F).var()
        
        # If blurry, use special enhancement
        if blur_value < 100:  # Threshold for blur detection
            return enhance_blurry_image(image)
        
        # Normal preprocessing for non-blurry images
        # Convert to grayscale
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
        else:
            gray = img
        
        # Denoising
        denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
        
        # Thresholding
        _, thresholded = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Sharpening
        kernel = np.array([[0, -1, 0], [-1, 5, -1], [0, -1, 0]])
        sharpened = cv2.filter2D(thresholded, -1, kernel)
        
        return Image.fromarray(sharpened)
    except Exception as e:
        st.error(f"Image preprocessing failed: {str(e)}")
        return image

def get_pdf_text(pdf_docs):
    """Extract text from a list of uploaded PDF documents."""
    raw_text = ""
    for pdf in pdf_docs:
        reader = PdfReader(pdf)
        for page in reader.pages:
            text = page.extract_text()
            if text:  # Make sure we're not adding None values
                raw_text += text
    return raw_text

def get_image_text(image_files):
    """Extract text from images with orientation correction"""
    raw_text = ""
    for image_file in image_files:
        try:
            if isinstance(image_file, bytes):
                # Handle camera capture
                image = Image.open(io.BytesIO(image_file))
            else:
                # Handle regular file upload
                image = Image.open(image_file)
            
            # Correct orientation and enhance
            corrected_image = correct_orientation(image)
            
            # Show original and processed images for comparison
            with st.expander("Document Processing Preview"):
                col1, col2 = st.columns(2)
                with col1:
                    st.image(corrected_image, caption="Original Document", use_column_width=True)
                
                processed_image = preprocess_image(corrected_image)
                with col2:
                    st.image(processed_image, caption="Processed Document", use_column_width=True)
            
            text = pytesseract.image_to_string(processed_image)
            if text:
                raw_text += text + "\n\n"
        except Exception as e:
            st.error(f"Error processing image: {str(e)}")
    return raw_text

def get_ppt_text(ppt_files):
    """Extract text from a list of uploaded PPT files."""
    raw_text = ""
    for ppt_file in ppt_files:
        presentation = Presentation(ppt_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw_text += shape.text + "\n"
    return raw_text

def get_video_text(video_files):
    """Extract frames from videos and analyze them"""
    raw_text = ""
    
    for video_file in video_files:
        try:
            # Save the uploaded file to a temporary location
            with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as temp_video:
                temp_video.write(video_file.getvalue())
                video_path = temp_video.name
            
            # Open video file
            cap = cv2.VideoCapture(video_path)
            frame_count = 0
            processed_frames = 0
            
            # Process at 1 frame per second
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_interval = int(fps)  # 1 frame per second
            
            while cap.isOpened():
                ret, frame = cap.read()
                if not ret:
                    break
                
                frame_count += 1
                
                # Only process every nth frame (1fps)
                if frame_count % frame_interval != 0:
                    continue
                
                # Convert frame to PIL Image
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                pil_image = Image.fromarray(frame_rgb)
                
                # Process the frame
                processed_image = preprocess_image(pil_image)
                
                # Extract text from frame
                frame_text = pytesseract.image_to_string(processed_image)
                if frame_text.strip():
                    raw_text += f"\nFrame {frame_count}:\n{frame_text}\n"
                    processed_frames += 1
                
                # Show progress
                if processed_frames % 5 == 0:
                    st.write(f"Processed {processed_frames} frames from {video_file.name}")
            
            cap.release()
            
            if not raw_text:
                raw_text = f"Could not extract text from video frames: {video_file.name}"
            else:
                raw_text = f"Extracted text from {processed_frames} frames:\n{raw_text}"
            
        except Exception as e:
            st.error(f"Error processing video {video_file.name}: {str(e)}")
            raw_text += f"\nError processing video: {video_file.name}"
        finally:
            # Clean up temporary files
            try:
                if os.path.exists(video_path):
                    os.unlink(video_path)
            except:
                pass
    
    return raw_text

def get_text_chunks(text):
    """Split text into manageable chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    """Create a vector store from text chunks."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embeddings)
    vector_store.save_local("faiss_index")
    return vector_store  # Return the vector store

def get_conversational_chain():
    """Create an advanced conversational chain with ICU/ER-specific capabilities and enhanced response generation"""
    
    prompt_template = """You are MediCare-ICU Pro, an advanced AI clinical assistant designed for ICU and Emergency Room (ER) scenarios. You provide **precise, actionable, and time-sensitive** medical intelligence ONLY for medical documents.  

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **DOCUMENT VERIFICATION PROTOCOL**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    1. First analyze if the provided context contains valid medical content:
       - Medical reports (labs, imaging, clinical notes)
       - Medication lists
       - Physician progress notes
       - Nursing documentation
       - Discharge summaries,medical bills 
       - Operative reports
       - Other recognized medical documentation
    
    2. If non-medical content is detected:
       - Financial documents (bank statements, bills)
       - Personal correspondence
       - Legal documents
       - Unrelated images/text
       - Commercial documents
       → IMMEDIATELY RESPOND:
           "⚠️ Non-Medical Document Detected: 
           The uploaded content appears to be non-medical in nature. 
           Please upload only medical reports for clinical analysis."

    CONTEXT FROM DOCUMENTS:
    {context}

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **MEDICAL CONTENT ANALYSIS**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    Only proceed if medical content is confirmed:

    USER QUERY:
    {question}

    USER TYPE: {user_type}  # (Healthcare Professional / Nurse / Patient / Family Member)

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **RESPONSE PROTOCOL v6.0**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

    **FOR DOCTORS (ICU/ER):**
    1. **Critical Risk Assessment**:
       - Highlight **life-threatening risks** (e.g., sepsis, cardiac arrest) with ⚠️.
       - Provide **clear thresholds** for lab values (e.g., lactate > 4 mmol/L, MAP < 65 mmHg).
       - Calculate and display relevant scores (e.g., SOFA, APACHE II, qSOFA) when applicable.
       - Suggest **immediate interventions** with time-sensitivity indicators.
       
    2. **Medication Management**:
       - Include **IV/IM medications** with precise dosing:
         - Drug name, dosage (weight-based if applicable), route, frequency, and duration.
         - Example: "IV Norepinephrine 0.1-2 mcg/kg/min, titrate to MAP > 65 mmHg."
       - Highlight **contraindications** and **drug interactions** with severity ratings.
       - Provide **alternative options** for allergies or contraindications.
       - Include renal/hepatic dose adjustments when appropriate.
       
    3. **Diagnostic Pathway**:
       - Summarize **key findings** from labs, imaging, and reports with clinical significance.
       - Suggest **next steps** in order of priority with timeframes.
       - Include **ICD-10 codes** for primary and secondary diagnoses.
       - Recommend further diagnostic tests with clinical rationale.
       
    4. **Emergency Protocols**:
       - Provide **step-by-step guidance** for critical scenarios with time-based actions.
       - Include **equipment preparation** and **team role assignment** recommendations.
       - Reference relevant guidelines (e.g., AHA, ACLS, ATLS) with protocol version.
       - Outline contingency plans for treatment failures.

    **FOR NURSES:**
    1. **Monitoring Parameters**:
       - Specify vital sign targets and frequency of monitoring.
       - Include early warning signs that require immediate physician notification.
       - Provide nursing-specific interventions and documentation requirements.
       - Suggest specific nursing diagnoses based on presented data.
       
    2. **Medication Administration**:
       - Highlight nursing considerations for medication administration.
       - Include compatibility information for IV medications.
       - Specify required monitoring during and after administration.
       - Note medications requiring special handling or precautions.

    **FOR PATIENTS/FAMILIES:**
    1. **Simplified Clinical Picture**:
       - Use **non-technical language** with analogies to explain complex conditions.
       - Explain rationale behind treatments in accessible terms.
       - Provide realistic expectations for recovery and timeline.
       - Address common concerns and questions specific to the condition.

    2. **Medication Guidance**:
       - Provide **basic instructions** focused on adherence and safety.
       - Explain purpose of medications in relation to symptoms.
       - Highlight important side effects to watch for.
       - Emphasize follow-up care requirements.

    3. **Supportive Communication**:
       - Offer **clear next steps** with emphasis on ongoing care.
       - Provide appropriate reassurance while maintaining accuracy.
       - Include resources for additional support and education.
       - Use empathetic language that acknowledges emotional impact.

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **CLINICAL DECISION SUPPORT MODULE**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    1. **Evidence Integration**:
       - Incorporate latest clinical guidelines and evidence-based recommendations.
       - Indicate level of evidence (e.g., Class I, IIa, IIb, III; Level A, B, C).
       - Highlight recent changes in clinical practice relevant to the case.
       - Calculate pre-test and post-test probabilities when applicable.
       
    2. **Differential Diagnosis**:
       - Rank differential diagnoses by likelihood based on presented data.
       - Include "can't miss" diagnoses even if less likely.
       - Specify key clinical features supporting or refuting each diagnosis.
       - Suggest targeted testing to narrow differentials efficiently.
       
    3. **Ventilator Management**:
       - Provide ventilator setting recommendations based on condition.
       - Include lung-protective strategies and weaning parameters.
       - Suggest arterial blood gas interpretation with ventilator adjustments.
       - Address ventilator-associated complications and prevention.
       
    4. **Fluid Management**:
       - Calculate fluid requirements based on weight, condition, and electrolyte status.
       - Recommend type, rate, and volume of fluids with monitoring parameters.
       - Address special considerations (heart failure, renal failure, etc.).
       - Provide triggers for reassessment and adjustment.

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **MEDICATION MODULE**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    1. **Critical Care Medications**:
       - Provide weight-based dosing with precision:
         - Example: "IV Furosemide 0.5-1 mg/kg (max 80mg), repeat in 6 hours if needed."
       - Include loading doses when applicable.
       - Specify renal/hepatic adjustments with GFR or Child-Pugh thresholds.
       - Include specific monitoring parameters with frequency.
       
    2. **Drug Interaction Analysis**:
       - Stratify interactions by severity (Mild ⚠️, Moderate ⚠️⚠️, Severe ⚠️⚠️⚠️).
       - Specify mechanism of interaction and clinical significance.
       - Recommend monitoring or management strategies for unavoidable interactions.
       - Include pharmacokinetic considerations for critically ill patients.
       
    3. **Sedation and Analgesia**:
       - Recommend goal-directed sedation with assessment scales (RASS, SAS).
       - Include sedation protocols with daily interruption criteria.
       - Provide pain management strategies specific to critical care.
       - Address delirium prevention and management.

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **RESPONSE STRUCTURE**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    1. **Critical Summary** (1-2 sentences highlighting most urgent findings/actions)
    2. **Clinical Assessment** (organized by system or priority)
    3. **Action Items** (numbered, in order of priority with timeframes)
    4. **Contingency Planning** (what-if scenarios and responses)
    5. **Follow-up Recommendations** (monitoring, reassessment times)

    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    **AI ETHICS & COMPLIANCE**
    ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    - Follow strict **medical ethics and confidentiality** (HIPAA & GDPR-compliant).
    - Ensure responses are **evidence-based** and aligned with current clinical guidelines.
    - Always recommend consulting a **licensed medical provider** for definitive care decisions.
    - Acknowledge limitations of AI-assisted analysis and emphasize human judgment.
    - Indicate when recommendations have significant practice variations or controversy.
    Generate either:
    1. A **Non-Medical Document Alert** if inappropriate content is detected, OR
    2. A **precise, ICU/ER-focused medical response** only if valid medical content is confirmed

    Generate a **precise, ICU/ER-focused response** based on the given context and user query:
    """
    
    
    model = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.2,
        top_k=30,
        top_p=0.95,
        max_output_tokens=2048
    )
    
    prompt = PromptTemplate(template=prompt_template, 
                          input_variables=["context", "question", "user_type"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def get_mini_chat_response(question):
    """Get response from AI medicoin model."""
    model = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.2,
        top_k=30,
        top_p=0.95,
        max_output_tokens=2048
    )
    response = model.invoke(question)
    return response.content

def camera_upload():
    """Handle camera upload with preview and confirmation"""
    col1, col2 = st.columns(2)
    if 'camera_on' not in st.session_state:
        st.session_state.camera_on = False

    with col1:
        if st.button("Turn Camera On"):
            st.session_state.camera_on = True
        if st.button("Turn Camera Off"):
            st.session_state.camera_on = False

    if st.session_state.camera_on:
        picture = st.camera_input("Take a picture of your document")
        
        if picture:
            with col2:
                st.write("Preview:")
                image = Image.open(io.BytesIO(picture.getvalue()))
                st.image(image, caption="Original Capture", width=200)
                
                if st.button("Confirm Upload"):
                    return picture
    return None

def set_medical_background():
    """Add medical-themed background and styles to the app."""
    with open("styles.css") as f:
        st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)

def main():
    # Page configuration with medical icon
    st.set_page_config(
        page_title="AI-Powered ICU & ER Assistant",
        page_icon="🏥",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Apply custom styles
    set_medical_background()
    
    # Handle password reset
    handle_password_reset()
    
    # Check authentication
    if not auth.user:
        st.title("MediCare-ICU Assistant")
        st.markdown("### AI-powered clinical decision support")
        show_auth()
        st.stop()
    
    # Add sign out button to sidebar
    with st.sidebar:
        if auth.user:
            st.markdown(f"**Welcome, {auth.user.get('firstName', 'User')}**")
            if st.button("Sign Out"):
                auth.sign_out()
    # Set medical-themed background
    set_medical_background()
    
    # Custom header with medical logo
    st.markdown("""
    <div class="main-header-container">
        <img src="https://cdn-icons-png.flaticon.com/512/2906/2906274.png" class="hospital-logo" alt="Hospital Logo">
        <h1 class="main-header">
            <span>🏥</span> MediCare-ICU Assistant
        </h1>
        <img src="https://cdn-icons-png.flaticon.com/512/2906/2906274.png" class="hospital-logo" alt="Hospital Logo">
    </div>
    <p class="sub-header">
        AI-powered clinical decision support for intensive care environments
    </p>
    """, unsafe_allow_html=True)
    
    # Modern disclaimer card
    st.markdown("""
    <div class="disclaimer">
        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.5rem;">
            <span style="font-size: 1.25rem;">⚠️</span>
            <strong style="font-size: 1rem;">MEDICAL DISCLAIMER</strong>
        </div>
        <p style="margin: 0;">
            This AI assistant provides informational support only and is not a substitute for professional medical judgment. 
            Always consult with qualified healthcare providers for diagnosis and treatment. 
            In emergencies, contact your local emergency services immediately.
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Create tabs with icons
    tabs = st.tabs([
        "📋 Chat Interface", 
        "💬 Mini Chat", 
        "📚 User Guide"
    ])
    
    with tabs[0]:
        # Two-column layout
        col1, col2 = st.columns([1, 2], gap="large")
        
        with col1:
            # Document management card
            with st.container():
                st.markdown("""
                <div class="card">
                    <h3 class="card-title">
                        <span>📁</span> Document Management
                    </h3>
                """, unsafe_allow_html=True)
                
                # User type selection
                user_type = st.radio(
                    "👤 I am a:",
                    ["Healthcare Professional", "Nurse", "Patient/Family"],
                    index=0,
                    help="Tailors responses to your clinical knowledge level"
                )
                
                # Camera upload section
                with st.expander("📷 Capture Document with Camera", expanded=False):
                    st.info("Position your document clearly in the camera view")
                    camera_capture = camera_upload()
                    if camera_capture:
                        st.markdown("""
                        <div class="success-banner">
                            <span>✅</span> Document captured successfully!
                        </div>
                        """, unsafe_allow_html=True)
                        if 'camera_docs' not in st.session_state:
                            st.session_state.camera_docs = []
                        st.session_state.camera_docs.append(camera_capture)
                
                # Display camera captures
                camera_files = []
                if 'camera_docs' in st.session_state and st.session_state.camera_docs:
                    st.write("**Camera Captures:**")
                    cols = st.columns(4)
                    for i, doc in enumerate(st.session_state.camera_docs):
                        with cols[i % 4]:
                            st.image(doc, caption=f"Capture {i+1}", width=100)
                            if st.button(f"❌ Remove {i+1}", key=f"remove_{i}"):
                                st.session_state.camera_docs.pop(i)
                                st.rerun()
                    camera_files = st.session_state.camera_docs.copy()
                
                # File uploader with modern styling
                uploaded_files = st.file_uploader(
                    "📤 Or upload documents directly",
                    type=["pdf", "png", "jpeg", "jpg", "ppt", "pptx", "mp4", "avi", "mov"],
                    accept_multiple_files=True,
                    help="Supported formats: PDF, Images, PowerPoint, Videos"
                )
                
                # Process documents button
                if st.button("🔍 Process Documents", use_container_width=True):
                    all_files = []
                    if uploaded_files:
                        all_files.extend(uploaded_files)
                    if camera_files:
                        all_files.extend(camera_files)
                    
                    if not all_files:
                        st.error("Please upload at least one document first.")
                    else:
                        with st.spinner("🔬 Analyzing medical documents..."):
                            try:
                                raw_text = ""
                                pdf_docs = [file for file in all_files if (hasattr(file, 'type') and file.type == "application/pdf")]
                                image_files = [file for file in all_files if (hasattr(file, 'type') and file.type in ["image/png", "image/jpeg", "image/jpg"]) or isinstance(file, bytes)]
                                ppt_files = [file for file in all_files if (hasattr(file, 'type') and file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"])]
                                video_files = [file for file in all_files if (hasattr(file, 'type') and file.type in ["video/mp4", "video/avi", "video/mov"])]
                                
                                if pdf_docs:
                                    raw_text += get_pdf_text(pdf_docs)
                                if image_files:
                                    raw_text += get_image_text(image_files)
                                if ppt_files:
                                    raw_text += get_ppt_text(ppt_files)
                                if video_files:
                                    raw_text += get_video_text(video_files)
                                
                                if not raw_text.strip():
                                    st.error("Could not extract text from the uploaded files.")
                                else:
                                    text_chunks = get_text_chunks(raw_text)
                                    st.session_state.vector_store = get_vector_store(text_chunks)  # Store in session state
                                    st.session_state.docs_processed = True
                                    st.markdown(f"""
                                    <div class="success-banner">
                                        <span>✅</span> {len(all_files)} document(s) processed successfully!
                                    </div>
                                    """, unsafe_allow_html=True)
                            except Exception as e:
                                st.error(f"An error occurred: {str(e)}")
                
                # Features card
                with st.container():
                    st.markdown("""
                    <div class="card">
                        <h3 class="card-title">
                            <span>💡</span> Features
                        </h3>
                        <div style="display: flex; flex-direction: column; gap: 0.75rem;">
                            <div style="display: flex; align-items: start; gap: 0.5rem;">
                                <span>🔍</span>
                                <div>
                                    <strong>Document Analysis</strong>
                                    <p style="margin: 0; font-size: 0.85rem; color: var(--light-text);">
                                        Extract insights from medical records, labs, and imaging
                                    </p>
                                </div>
                            </div>
                            <div style="display: flex; align-items: start; gap: 0.5rem;">
                                <span>💊</span>
                                <div>
                                    <strong>Medication Support</strong>
                                    <p style="margin: 0; font-size: 0.85rem; color: var(--light-text);">
                                        Identify interactions, dosages, and alternatives
                                    </p>
                                </div>
                            </div>
                            <div style="display: flex; align-items: start; gap: 0.5rem;">
                                <span>🩺</span>
                                <div>
                                    <strong>Clinical Decision Support</strong>
                                    <p style="margin: 0; font-size: 0.85rem; color: var(--light-text);">
                                        Evidence-based recommendations and protocols
                                    </p>
                                </div>
                            </div>
                            <div style="display: flex; align-items: start; gap: 0.5rem;">
                                <span>🎤</span>
                                <div>
                                    <strong>Voice Assistant</strong>
                                    <p style="margin: 0; font-size: 0.85rem; color: var(--light-text);">
                                        Speak your questions and get voice responses
                                    </p>
                                </div>
                            </div>
                            <div style="display: flex; align-items: start; gap: 0.5rem;">
                                <span>📷</span>
                                <div>
                                    <strong>Image Enhancement</strong>
                                    <p style="margin: 0; font-size: 0.85rem; color: var(--light-text);">
                                        Automatically enhance blurry document images
                                    </p>
                                </div>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                
                # Quick prompts
                if st.session_state.get('docs_processed', False):
                    with st.container():
                        st.markdown("""
                        <div class="card">
                            <h3 class="card-title">
                                <span>🚀</span> Quick Prompts
                        """, unsafe_allow_html=True)
                        
                        if user_type == "Healthcare Professional":
                            cols = st.columns(2)
                            with cols[0]:
                                if st.button("📈 Lab Trends", use_container_width=True):
                                    st.session_state.preset_question = "Summarize the most recent lab results and highlight any abnormal values."
                            with cols[1]:
                                if st.button("💊 Med Conflicts", use_container_width=True):
                                    st.session_state.preset_question = "Are there any potential conflicts between the current medications?"
                            cols = st.columns(2)
                            with cols[0]:
                                if st.button("📋 Treatment Plan", use_container_width=True):
                                    st.session_state.preset_question = "What is the recommended treatment protocol for this condition?"
                            with cols[1]:
                                if st.button("⚠️ Risk Factors", use_container_width=True):
                                    st.session_state.preset_question = "What are the key risk factors to monitor in this case?"
                        elif user_type == "Nurse":
                            cols = st.columns(2)
                            with cols[0]:
                                if st.button("🩺 Patient Care", use_container_width=True):
                                    st.session_state.preset_question = "What are the best practices for patient care in this scenario?"
                            with cols[1]:
                                if st.button("💊 Medication Administration", use_container_width=True):
                                    st.session_state.preset_question = "What are the guidelines for administering these medications?"
                            cols = st.columns(2)
                            with cols[0]:
                                if st.button("📋 Care Plan", use_container_width=True):
                                    st.session_state.preset_question = "What is the recommended care plan for this patient?"
                            with cols[1]:
                                if st.button("⚠️ Safety Protocols", use_container_width=True):
                                    st.session_state.preset_question = "What safety protocols should be followed in this case?"
                        else:
                            cols = st.columns(2)
                            with cols[0]:
                                if st.button("📝 Explain Diagnosis", use_container_width=True):
                                    st.session_state.preset_question = "Can you explain what the diagnosis means in simple terms?"
                            with cols[1]:
                                if st.button("💊 Med Side Effects", use_container_width=True):
                                    st.session_state.preset_question = "What are the common side effects of the prescribed medications?"
                            cols = st.columns(2)
                            with cols[0]:
                                if st.button("🔄 Recovery Process", use_container_width=True):
                                    st.session_state.preset_question = "What should we expect during the recovery process?"
                            with cols[1]:
                                if st.button("❓ Questions for Doctor", use_container_width=True):
                                    st.session_state.preset_question = "What questions should I ask my doctor at the next appointment?"
                        
                        st.markdown("</div>", unsafe_allow_html=True)

        with col2:
            # Chat interface
            st.markdown("""
            <div class="card">
                <h3 class="card-title">
                    <span>💬</span> Medical Consultation
            """, unsafe_allow_html=True)
            
            # Initialize chat history
            if 'messages' not in st.session_state:
                st.session_state.messages = []
            
            # Initialize docs_processed state
            if 'docs_processed' not in st.session_state:
                st.session_state.docs_processed = False
            
            # Initialize preset question state
            if 'preset_question' not in st.session_state:
                st.session_state.preset_question = ""
            
            # Initialize vector store in session state if not exists
            if 'vector_store' not in st.session_state:
                st.session_state.vector_store = None
            
            # Add voice assistant toggle
            st.sidebar.markdown("### Voice Assistant")
            use_voice_assistant = st.sidebar.checkbox("Enable Voice Assistant", value=True)
            
            # Voice assistant button
            if use_voice_assistant and st.session_state.docs_processed:
                if st.sidebar.button("🎤 Speak Your Question", key="voice_input_btn"):
                    voice_input = get_voice_input()
                    if voice_input:
                        st.session_state.preset_question = voice_input
                        st.rerun()

            # Display chat history with modern styling
            for message in st.session_state.messages:
                if message['role'] == 'user':
                    st.markdown(f"""
                    <div class="chat-user">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.25rem;">
                            <span style="font-size: 1rem;">👤</span>
                            <strong>You</strong>
                        </div>
                        {message['content']}
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                    <div class="chat-assistant">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.25rem;">
                            <span style="font-size: 1rem;">🏥</span>
                            <strong>MediCare-ICU</strong>
                        </div>
                        {message['content']}
                    </div>
                    """, unsafe_allow_html=True)
            
            # Chat input with conditional disabling
            if st.session_state.preset_question:
                user_question = st.session_state.preset_question
                st.session_state.preset_question = ""  # Clear after use
            else:
                user_question = st.chat_input(
                    "Type your medical question here...",
                    disabled=not st.session_state.docs_processed,
                    key="main_input"
                )
            
            # Process query and display response
            if user_question:
                # Add user message to chat
                st.session_state.messages.append({"role": "user", "content": user_question})
                
                with st.spinner("🔍 Analyzing medical information..."):
                    try:
                        if st.session_state.vector_store is None:
                            st.error("Please process documents first")
                        else:
                            # Get embeddings and search for relevant documents
                            docs = st.session_state.vector_store.similarity_search(user_question, k=5)
                            
                            # Get the conversational chain and generate response
                            chain = get_conversational_chain()
                            response = chain(
                                {
                                    "input_documents": docs, 
                                    "question": user_question,
                                    "user_type": user_type
                                },
                                return_only_outputs=True
                            )
                            
                            response_text = response["output_text"]
                            
                            # Add assistant response to chat history
                            st.session_state.messages.append({"role": "assistant", "content": response_text})
                            
                            # Display the response text
                            st.markdown(f"""
                            <div class="chat-assistant">
                                <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.25rem;">
                                    <span style="font-size: 1rem;">🏥</span>
                                    <strong>MediCare-ICU</strong>
                                </div>
                                {response_text}
                            </div>
                            """, unsafe_allow_html=True)
                            
                            # Generate and save audio file
                            audio_filename = f"response_{datetime.now().strftime('%Y%m%d_%H%M%S')}.mp3"
                            saved_file = generate_audio_file(response_text, audio_filename)

                            if saved_file:
                                autoplay_audio(saved_file)
                            
                            # Emergency alert
                            if "⚠️ URGENT:" in response_text:
                                st.markdown("""
                                <div class="emergency">
                                    <div style="display: flex; align-items: center; gap: 0.5rem;">
                                        <span style="font-size: 1.25rem;">🚨</span>
                                        <strong>EMERGENCY ALERT</strong>
                                    </div>
                                    <p style="margin: 0.5rem 0 0 0;">
                                        This response indicates a potential medical emergency. 
                                        Please seek immediate medical attention.
                                    </p>
                                </div>
                                """, unsafe_allow_html=True)
                        
                        # Force rerun to update the chat display
                        st.rerun()
                    
                    except Exception as e:
                        st.error(f"An error occurred: {str(e)}")
            
            else:
                if not st.session_state.docs_processed:
                    st.info("ℹ️ Please upload and process medical documents first")
            
            # Follow-up section
            st.markdown("""
            <div style="margin-top: 1.5rem;">
                <h4 style="color: var(--primary); margin-bottom: 0.5rem;">
                    <span>🔄</span> Follow-up Questions
                </h4>
            """, unsafe_allow_html=True)
            
            follow_up_question = st.chat_input(
                "Ask a follow-up question...",
                key="follow_up_input"
            )
            
            if follow_up_question:
                st.session_state.messages.append({"role": "user", "content": follow_up_question})
                
                with st.spinner("💭 Generating response..."):
                    try:
                        if st.session_state.vector_store is None:
                            st.error("Please process documents first")
                        else:
                            docs = st.session_state.vector_store.similarity_search(follow_up_question, k=5)
                            chain = get_conversational_chain()
                            response = chain(
                                {
                                    "input_documents": docs, 
                                    "question": follow_up_question,
                                    "user_type": user_type
                                },
                                return_only_outputs=True
                            )
                            follow_up_response = response["output_text"]
                            st.session_state.messages.append({"role": "assistant", "content": follow_up_response})
                        
                        # Force rerun to update the chat display
                        st.rerun()
                    except Exception as e:
                        st.error(f"Error: {str(e)}")
            
            st.markdown("</div></div>", unsafe_allow_html=True)

    with tabs[1]:
        # Mini chat interface
        st.markdown("""
        <div class="card">
            <h3 class="card-title">
                <span>💬</span> Quick Medical Questions
            </h3>
            <p style="color: var (--light-text); margin-top: -0.5rem; margin-bottom: 1rem;">
                Ask general medical questions without document processing
            </p>
        """, unsafe_allow_html=True)
        
        # Initialize mini chat history
        if 'mini_messages' not in st.session_state:
            st.session_state.mini_messages = []
        
        # Display mini chat history
        for message in st.session_state.mini_messages:
            if message['role'] == 'user':
                st.markdown(f"""
                <div class="chat-user">
                    <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.25rem;">
                        <span style="font-size: 1rem;">👤</span>
                        <strong>You</strong>
                    </div>
                    {message['content']}
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="mini-chat">
                    <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.25rem;">
                        <span style="font-size: 1rem;">💡</span>
                        <strong>MEDICOIN</strong>
                    </div>
                    {message['content']}
                </div>
                """, unsafe_allow_html=True)
        
        # Mini chat input
        mini_question = st.chat_input("Ask a general medical question...", key="mini_input")
        
        if mini_question:
            st.session_state.mini_messages.append({"role": "user", "content": mini_question})
            
            with st.spinner("💭 Generating response..."):
                try:
                    mini_response = get_mini_chat_response(mini_question)
                    st.session_state.mini_messages.append({"role": "assistant", "content": mini_response})
                    st.rerun()
                except Exception as e:
                    st.error(f"Error: {str(e)}")
        
        st.markdown("</div>", unsafe_allow_html=True)
        with tabs[2]:
        # User Guide with professional styling
            st.markdown("""
        <div class="user-guide-section">
            <h2 class="user-guide-title">
                <span>📚</span> MediCare-ICU User Guide
            </h2>
            <p style="color: var(--light-text); margin-bottom: 1.5rem;">
                A comprehensive guide to using the MediCare-ICU Assistant for optimal clinical decision support.
            </p>
            
<p>
MediCare-ICU Assistant - User Guide

Here's a clear, point-by-point version of the MediCare-ICU User Guide:

**MediCare-ICU Assistant - User Guide**

1. **Uploading Documents**
   - Supported formats: PDFs, images (PNG/JPG), PowerPoint files, videos
   - Two upload methods: File uploader or camera capture
   - *Pro Tip:* For images, ensure good lighting and focus for best OCR results

2. **Processing Documents**
   - Click "Process Documents" to analyze files
   - System extracts text and identifies key medical information
   - Tailored outputs:
     - *Healthcare Professionals:* Detailed clinical data (labs, meds, diagnoses)
     - *Patients/Families:* Simplified explanations of medical content

3. **Asking Questions**
   - Ask natural language questions after processing
   - Example queries:
     - "List current medications"
     - "Explain this diagnosis simply"
     - "Show potential drug interactions"
     - "Recommend treatment options"

4. **Understanding Responses**
   - Responses include:
     - Clinical insights with document references
     - Medication details (dosages, interactions, side effects)
     - Evidence-based treatment recommendations
     - Emergency alerts for critical findings
   - *Important:* Always verify critical findings with a healthcare provider

5. **Follow-up Questions**
   - System maintains conversation context
   - Ask clarifying questions like:
     - "Side effects of [medication]?"
     - "Alternative treatment options?"
   - Use Quick Prompts for common questions about labs/meds/treatments

**Troubleshooting**
- Document issues: Ensure PDFs have selectable text, use clear images
- Connection problems: Check internet/API key
- Performance: Large files may take longer to process

**Support**
- Technical: balubejagam@gmail.com | +91 7993650197
- Clinical: clinical@medicareicu.com | +91 9676712990


</p>
        </div>
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()

def create_reset_table():
    """Create password reset table."""
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS password_resets (
            email TEXT PRIMARY KEY,
            token TEXT NOT NULL,
            expires_at TIMESTAMP NOT NULL
        )
    """)
    conn.commit()
    conn.close()

create_reset_table()

def create_verification_table():
    """Create table for email verification tokens."""
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS email_verifications (
            email TEXT PRIMARY KEY,
            token TEXT NOT NULL,
            expires_at TIMESTAMP NOT NULL,
            verified BOOLEAN DEFAULT FALSE
        )
    """)
    conn.commit()
    conn.close()

create_verification_table()

def send_verification_email(email: str, token: str):
    """Send verification email with SMTP connection handling."""
    try:
        msg = MIMEMultipart()
        msg['From'] = EMAIL_USER
        msg['To'] = email
        msg['Subject'] = "Verify Your MediCare-ICU Account"
        
        verification_link = f"{BASE_URL}/?verify_token={token}"
        body = f"""
        <h2>Account Verification</h2>
        <p>Please click the link below to verify your email:</p>
        <p><a href="{verification_link}">Verify Email</a></p>
        <p>This link will expire in 24 hours.</p>
        """
        
        msg.attach(MIMEText(body, 'html'))
        
        with smtplib.SMTP(EMAIL_HOST, EMAIL_PORT, timeout=10) as server:
            server.ehlo()
            server.starttls()
            server.login(EMAIL_USER, EMAIL_PASS)
            server.send_message(msg)
        return True
    except Exception as e:
        st.error(f"Failed to send verification email: {str(e)}")
        return False

def handle_email_verification():
    """Check for verification token in URL."""
    query_params = st.query_params  # Updated from st.experimental_get_query_params
    if 'verify_token' in query_params:
        token = query_params['verify_token']
        
        conn = sqlite3.connect('medicare_users.db')
        c = conn.cursor()
        c.execute("""
            SELECT email FROM email_verifications 
            WHERE token=? AND expires_at > datetime('now') AND verified=FALSE
        """, (token,))
        result = c.fetchone()
        
        if result:
            email = result[0]
            c.execute("""
                UPDATE email_verifications 
                SET verified=TRUE 
                WHERE email=?
            """, (email,))
            c.execute("""
                UPDATE users 
                SET verified=TRUE 
                WHERE email=?
            """, (email,))
            conn.commit()
            st.success("Email verified successfully! You can now login.")
            st.query_params.clear()  # Updated from st.experimental_set_query_params
            st.rerun()
        else:
            st.error("Invalid or expired verification link")
        conn.close()

# Update sign-up function to include email verification
def sign_up(self, username, password, email, first_name, last_name, role="patient"):
    if not is_valid_email(email):
        return False, "Invalid email format"
    
    conn = sqlite3.connect('medicare_users.db')
    c = conn.cursor()
    try:
        c.execute("""
            INSERT INTO users 
            (username, password, email, first_name, last_name, role, created_at, verified) 
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (username, hash_password(password), email, first_name, last_name, role, datetime.now(), False))
        
        token = generate_reset_token()
        expires_at = datetime.now() + timedelta(hours=24)
        c.execute("""
            INSERT INTO email_verifications 
            (email, token, expires_at) 
            VALUES (?, ?, ?)
        """, (email, token, expires_at))
        
        conn.commit()
        
        if send_verification_email(email, token):
            return True, "Account created! Please check your email for verification."
        else:
            return True, "Account created, but failed to send verification email. Please contact support."
    except sqlite3.IntegrityError:
        return False, "Username or email already exists"
    finally:
        conn.close()

# Update authentication UI to include verification status
def show_auth():
    tabs = st.tabs(["Login", "Sign Up", "Resend Verification"])
    
    with tabs[0]:  # Login
        with st.form("login_form"):
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            
            if st.form_submit_button("Login"):
                conn = sqlite3.connect('medicare_users.db')
                c = conn.cursor()
                c.execute("""
                    SELECT password, verified FROM users WHERE username=?
                """, (username,))
                result = c.fetchone()
                
                if result and verify_password(password, result[0]):
                    if result[1]:  # Verified
                        st.success("Logged in successfully!")
                        st.rerun()
                    else:
                        st.error("Please verify your email first")
                else:
                    st.error("Invalid credentials")
                conn.close()
    
    with tabs[1]:  # Sign Up
        # ...existing sign-up form code...
        pass
    
    with tabs[2]:  # Resend Verification
        with st.form("resend_form"):
            email = st.text_input("Your registered email")
            if st.form_submit_button("Resend Verification Email"):
                conn = sqlite3.connect('medicare_users.db')
                c = conn.cursor()
                c.execute("""
                    SELECT verified FROM users WHERE email=?
                """, (email,))
                result = c.fetchone()
                
                if result:
                    if result[0]:
                        st.warning("Email already verified")
                    else:
                        token = generate_reset_token()
                        expires_at = datetime.now() + timedelta(hours=24)
                        c.execute("""
                            UPDATE email_verifications 
                            SET token=?, expires_at=?, verified=FALSE
                            WHERE email=?
                        """, (token, expires_at, email))
                        conn.commit()
                        
                        if send_verification_email(email, token):
                            st.success("Verification email resent!")
                        else:
                            st.error("Failed to resend email")
                else:
                    st.error("Email not found")
                conn.close()

# Call handle_email_verification at the start of main()
handle_email_verification()
# ...existing code